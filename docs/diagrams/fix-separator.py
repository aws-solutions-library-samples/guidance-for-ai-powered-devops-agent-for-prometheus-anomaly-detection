#!/usr/bin/env python3
"""Post-process the Guidance PPTX so the AWS RA Validator can parse it.

Why this exists
---------------
The official guidance-architecture-diagram-template draws the horizontal separator
under the title as a CONNECTOR (<p:cxnSp>, reported as shape type LINE/9).
pptxgenjs can only emit lines as autoshapes (<p:sp> with prstGeom="line", reported as
AUTO_SHAPE/1). The RA Validator looks for a LINE-type shape spanning >=70% of the slide
width, so with a pptxgenjs line it reports:

    Separator Line  -> Failed
    Font Title      -> "Title starting with 'Guidance for': Not Found"

The title failure is a knock-on effect: the validator uses the separator to work out
which shapes are the title/description, so with no separator it cannot find the title.

This script swaps that one autoshape for a real straight connector, in place.

Usage: python3 fix-separator.py <pptx>
"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.dml.color import RGBColor

SEP_Y      = 1.30          # inches, must match the generator
SEP_X      = 0.12
SEP_W      = 9.60          # 72% of 13.333" (validator needs >= 70%); ends 9.72 < sidebar 9.82
SEP_COLOR  = RGBColor(0x54, 0x5B, 0x64)
SEP_WIDTH  = Pt(2)


def main(path):
    prs = Presentation(path)
    slide = prs.slides[0]
    sw_in = prs.slide_width / 914400

    # 1) drop the pptxgenjs line autoshape that acts as the separator
    removed = 0
    for sh in list(slide.shapes):
        if sh.width and sh.height == 0 and abs((sh.top or 0) / 914400 - SEP_Y) < 0.05:
            sh._element.getparent().remove(sh._element)
            removed += 1

    # 2) add a real connector in its place
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(SEP_X), Inches(SEP_Y), Inches(SEP_X + SEP_W), Inches(SEP_Y)
    )
    conn.name = "Separator"
    conn.line.color.rgb = SEP_COLOR
    conn.line.width = SEP_WIDTH

    prs.save(path)

    # 3) verify
    chk = Presentation(path)
    lines = [sh for sh in chk.slides[0].shapes
             if sh._element.tag.endswith('cxnSp')]
    ok = False
    for sh in lines:
        pct = sh.width / chk.slide_width * 100
        print(f"  separator: {sh.name} type={sh.shape_type} "
              f"x={sh.left/914400:.2f} w={sh.width/914400:.2f} ({pct:.1f}% of slide)")
        if pct >= 70:
            ok = True
    title = [sh.text for sh in chk.slides[0].shapes
             if sh.has_text_frame and sh.text.startswith("Guidance for")]
    print(f"  removed {removed} line autoshape(s); connector >=70% width: {ok}")
    print(f"  title above separator: {'found' if title else 'MISSING'}")
    return 0 if (ok and title) else 1


if __name__ == "__main__":
    sys.exit(main(sys.argv[1] if len(sys.argv) > 1 else "5g-rcf-architecture-guidance.pptx"))
